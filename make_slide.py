# -*- coding: utf-8 -*-
"""Sinh slide PowerPoint (.pptx) DAY DU cho de tai Lap trinh nhan Linux.
   - Phu ca 4 phan, rieng Phan 4 trinh bay chuyen sau (7 slide)
   - Cac slide demo co KHUNG placeholder de dan anh chup man hinh."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE   = RGBColor(0x1F, 0x4E, 0x79)
LBLUE  = RGBColor(0xBF, 0xD7, 0xFF)
DARK   = RGBColor(0x22, 0x22, 0x22)
GRAY   = RGBColor(0x55, 0x55, 0x55)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x99, 0x55, 0x00)
CODEBG = RGBColor(0x2B, 0x2B, 0x2B)
CODEFG = RGBColor(0xD4, 0xD4, 0xD4)

blank = prs.slide_layouts[6]

def _title_bar(s, title):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                             prs.slide_width, Inches(1.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.85))
    tp = tb.text_frame; tp.word_wrap = True
    p = tp.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(27); r.font.bold = True; r.font.color.rgb = WHITE

def _bullets(s, bullets, x, y, w, h):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.word_wrap = True
    for i, (txt, lvl) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        r = p.add_run()
        r.text = ('• ' if lvl == 0 else '– ') + txt
        r.font.size = Pt(19 - lvl*2)
        r.font.color.rgb = DARK if lvl == 0 else GRAY
        p.space_after = Pt(7)

def _code(s, code, x, y, w, h):
    cb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    cb.fill.solid(); cb.fill.fore_color.rgb = CODEBG; cb.line.color.rgb = GRAY
    tf = cb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.1)
    for i, ln in enumerate(code.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = ln if ln else ' '
        r.font.name = 'Consolas'; r.font.size = Pt(11); r.font.color.rgb = CODEFG

def _imgbox(s, hint, x, y, w, h):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
    box.line.color.rgb = RGBColor(0x99, 0x99, 0x99); box.line.width = Pt(1.5)
    box.line.dash_style = None
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = '🖼  CHÈN ẢNH DEMO'
    r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = ACCENT
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = hint
    r2.font.size = Pt(12); r2.font.italic = True; r2.font.color.rgb = GRAY

def content(title, bullets, code=None, img=None):
    s = prs.slides.add_slide(blank)
    _title_bar(s, title)
    if code or img:
        _bullets(s, bullets, Inches(0.55), Inches(1.45), Inches(7.0), Inches(5.6))
        if code:
            _code(s, code, Inches(7.75), Inches(1.45), Inches(5.15), Inches(5.5))
        else:
            _imgbox(s, img, Inches(7.75), Inches(2.1), Inches(5.15), Inches(4.2))
    else:
        _bullets(s, bullets, Inches(0.7), Inches(1.5), Inches(12.0), Inches(5.6))
    return s

def section(num, title):
    s = prs.slides.add_slide(blank)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.7),
                             prs.slide_width, Inches(2.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.95), Inches(11.7), Inches(1.7))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = num
    r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = LBLUE
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = title
    r2.font.size = Pt(32); r2.font.bold = True; r2.font.color.rgb = WHITE
    return s

# ===================== 1. BIA =====================
s = prs.slides.add_slide(blank)
box = s.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(11.7), Inches(2.2))
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = 'LẬP TRÌNH NHÂN LINUX'
r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = BLUE
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = ('Quản lý hệ thống (Shell) · Tiến trình – File – Socket – Network (C) · '
           'Module nhân · ARM Linux mô phỏng (Buildroot & QEMU)')
r2.font.size = Pt(18); r2.font.color.rgb = GRAY
box2 = s.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11.7), Inches(2.0))
tf2 = box2.text_frame; tf2.word_wrap = True
for i, ln in enumerate(['Môn học: Lập trình nhân Linux',
                        'GVHD: ...............    SVTH: ...............    Lớp: ...............',
                        'Hà Nội, 2026']):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = ln; r.font.size = Pt(15); r.font.color.rgb = DARK

# ===================== 2. NOI DUNG =====================
content('Nội dung trình bày', [
    ('Giới thiệu chung về Linux và hệ thống nhúng', 0),
    ('Phần 1 — Lập trình Shell quản lý hệ thống', 0),
    ('Phần 2 — Lập trình C: tiến trình, file, socket, network', 0),
    ('Phần 3 — Xây dựng module nhân Linux', 0),
    ('Phần 4 — Hệ thống ARM Linux mô phỏng (Buildroot & QEMU)', 0),
    ('Kết luận', 0),
])

# ===================== 3. GIOI THIEU =====================
content('Giới thiệu chung', [
    ('Linux: hệ điều hành mã nguồn mở, miễn phí, bảo mật, tùy biến cao', 0),
    ('Cấu trúc 3 lớp: Kernel – Shell – Application', 1),
    ('Hiện diện: máy chủ, cloud, di động (Android), thiết bị nhúng', 1),
    ('Hệ thống nhúng (Embedded Linux):', 0),
    ('Tinh gọn, chiếm ít tài nguyên, độc lập nền tảng', 1),
    ('Dùng trong điện thoại, router, IoT, ô tô…', 1),
    ('Mục tiêu đề tài: Shell → C (system call) → Kernel → HĐH nhúng ARM', 0),
])

# ===================== PHAN 1 =====================
section('PHẦN 1', 'Lập trình Shell quản lý hệ thống')
content('Phần 1 — Lập trình Shell', [
    ('Quản lý file: tạo, đổi tên, sao chép, xóa, phân quyền, nén', 0),
    ('Lập lịch tác vụ với crontab', 0),
    ('Thiết lập thời gian hệ thống & đồng bộ NTP', 0),
    ('Cài đặt/gỡ bỏ chương trình tự động (apt, dpkg)', 0),
    ('Tổ chức dạng menu; menu.sh gọi 4 script con', 0),
], code='''# them tac vu crontab
them() {
  read -p "Lich: " lich
  (crontab -l 2>/dev/null;
   echo "$lich") | crontab -
}
# cai dat goi
cai_dat() {
  sudo apt-get install -y "$goi"
}''')
content('Phần 1 — Kết quả thực thi', [
    ('4 chương trình chạy theo menu, thao tác đúng trên hệ thống', 0),
    ('In thông báo [OK] / [Lỗi] sau mỗi thao tác', 0),
    ('Cần quyền sudo cho đổi giờ và cài/gỡ gói', 1),
], img='Chụp terminal chạy menu.sh và một thao tác (vd: quản lý file hoặc crontab)')

# ===================== PHAN 2 =====================
section('PHẦN 2', 'Lập trình C: Tiến trình – File – Socket – Network')
content('Phần 2 — Quản lý tiến trình & file', [
    ('Quản lý file qua system call: stat, chmod, opendir/readdir', 0),
    ('Quản lý tiến trình:', 0),
    ('Liệt kê tiến trình (ps)', 1),
    ('Tạo tiến trình con: fork() + execl()', 1),
    ('Gửi tín hiệu / tiêu diệt: kill(); chờ: waitpid()', 1),
], code='''pid_t pid = fork();
if (pid == 0) {        /* con */
  execl("/bin/sh","sh",
        "-c", lenh, NULL);
} else {               /* cha */
  int st;
  waitpid(pid, &st, 0);
}''')
content('Phần 2 — Socket & Network', [
    ('Network: liệt kê giao diện & IP bằng getifaddrs()', 0),
    ('Socket TCP — mô hình Client/Server:', 0),
    ('Server: socket → bind → listen → accept', 1),
    ('Client: socket → connect', 1),
    ('Trao đổi tin nhắn 2 chiều; gõ "bye" để kết thúc', 1),
    ('Biên dịch bằng Makefile (make / make clean)', 0),
], code='''/* server */
bind(fd,&addr,len);
listen(fd, 3);
accept(fd, ...);

/* client */
connect(s,&serv,len);''')
content('Phần 2 — Kết quả thực thi', [
    ('5 chương trình C biên dịch & chạy thành công', 0),
    ('Demo socket: server và client nhắn tin qua lại', 0),
], img='Chụp 2 cửa sổ terminal: ./server và ./client đang nhắn tin')

# ===================== PHAN 3 =====================
section('PHẦN 3', 'Xây dựng Module nhân Linux')
content('Phần 3 — Module nhân Linux', [
    ('Module nhân (LKM): nạp/gỡ khi hệ thống đang chạy', 0),
    ('Vòng đời: module_init (__init) ↔ module_exit (__exit)', 0),
    ('Nhận tham số đầu vào: module_param', 0),
    ('Ghi log bằng printk → xem bằng dmesg', 0),
    ('Demo: tính tổng 1..n và liệt kê số nguyên tố trong kernel', 0),
], code='''module_param(n,int,0644);

__init hello_init(void){
 printk("n=%d\\n", n);
 /* tinh toan ... */
 return 0;
}
module_init(hello_init);
module_exit(hello_exit);''')
content('Phần 3 — Biên dịch & Kết quả', [
    ('make → hello_module.ko', 0),
    ('sudo insmod hello_module.ko n=10 ten="Linux"', 0),
    ('dmesg | tail  → xem log của module', 0),
    ('sudo rmmod hello_module → gỡ module', 0),
    ('Log mẫu: Tong 1..10 = 55 | Nguyen to: 2 3 5 7', 1),
], img='Chụp terminal: insmod + dmesg hiển thị log, và dòng "Tam biet" khi rmmod')

# ===================== PHAN 4 (CHUYEN SAU - 7 SLIDE) =====================
section('PHẦN 4', 'Hệ thống ARM Linux mô phỏng bằng Buildroot & QEMU')

# 4.1
content('Phần 4 (1/7) — Tổng quan & kiến trúc', [
    ('Mục tiêu: dựng HĐH Linux nhúng cho ARM, chạy hoàn toàn bằng phần mềm', 0),
    ('Không cần phần cứng thật → phát triển & kiểm thử nhanh', 1),
    ('Kiến trúc đích: ARM 64-bit (ARMv8-A), CPU Cortex-A53', 0),
    ('Phần cứng GPIO giả lập: ARM PrimeCell PL061 @ 0x09030000 (8 đường)', 0),
    ('Hai công cụ chính: Buildroot (dựng image) + QEMU (giả lập)', 0),
], img='Sơ đồ tổng quan: Host (Buildroot) → image → QEMU chạy ARM Linux')

# 4.2
content('Phần 4 (2/7) — Buildroot', [
    ('Công cụ tự động hóa xây dựng hệ thống nhúng từ mã nguồn mở', 0),
    ('Tự tải mã nguồn, biên dịch chéo (cross-compile), đóng gói', 0),
    ('Tính năng: linh hoạt, tối ưu kích thước, tự động hóa, đa nền tảng', 0),
    ('Thành phần: Makefile, Config.in, .config (menuconfig)', 0),
    ('Đầu ra: Kernel + Root filesystem + Bootloader + phụ trợ', 0),
], img='Ảnh trang chủ/cấu trúc Buildroot hoặc cây thư mục output/')

# 4.3
content('Phần 4 (3/7) — QEMU', [
    ('Quick Emulator — phần mềm ảo hóa/giả lập mã nguồn mở', 0),
    ('Giả lập nhiều kiến trúc CPU: x86, ARM, MIPS, PowerPC…', 0),
    ('Hai thành phần: emulator (giả lập phần cứng) + hypervisor', 0),
    ('Khởi động: qemu-system-aarch64 -M virt -cpu cortex-a53 -nographic', 0),
    ('User-mode networking: máy chủ truy cập qua 10.0.2.2', 0),
    ('Thoát QEMU: Ctrl-a rồi x', 1),
], code='''cd output/images
./start-qemu.sh

# login: root
# (khong can mat khau)

# thoat: Ctrl-a x''')

# 4.4
content('Phần 4 (4/7) — Xây dựng image Linux', [
    ('B1: Tải Buildroot + cấu hình mặc định cho QEMU ARM64', 0),
    ('B2: menuconfig — thêm gói GPIO:', 0),
    ('Kernel > Linux Kernel Tools > gpio', 1),
    ('Libraries > Hardware handling > libgpiod (+ tools)', 1),
    ('Networking applications > openssh', 1),
    ('B3: linux-menuconfig — bật driver PL061', 0),
    ('B4: make (~30 phút) → B5: ./start-qemu.sh', 0),
], code='''make qemu_aarch64_virt\\
     _defconfig
make menuconfig
make linux-menuconfig
#  GPIO Support >
#  PrimeCell PL061
time make''')

# 4.5
content('Phần 4 (5/7) — Ứng dụng & biên dịch chéo', [
    ('Viết ứng dụng người dùng hello.c trên máy chủ', 0),
    ('Biên dịch bằng cross-compiler do Buildroot tạo (aarch64-…-gcc)', 0),
    ('Chuyển tệp vào QEMU: ssh/scp hoặc wget qua HTTP (10.0.2.2)', 0),
    ('Cấp quyền chmod +x rồi chạy trên máy ARM khách', 0),
], code='''aarch64-linux-gcc \\
   -o hello hello.c

# trong QEMU:
wget http://10.0.2.2:\\
     8080/hello
chmod +x hello
./hello   # Hello QEMU''')

# 4.6
content('Phần 4 (6/7) — Giao tiếp GPIO (libgpiod)', [
    ('gpiodetect → gpiochip0 [9030000.pl061] (8 lines)', 0),
    ('gpioinfo → 8 chân: unused, input, active-high', 0),
    ('blink.c: dùng libgpiod nhấp nháy chân 3, chu kỳ 200ms', 0),
    ('pl061-change.c: mmap /dev/mem @0x09030000, giám sát GPIODATA', 0),
    ('Kết quả: GPIODATA = 8 ↔ 0 (8 = 2³ của chân số 3)', 1),
], img='Chụp terminal QEMU: gpiodetect/gpioinfo và log GPIODATA=8/0 luân phiên')

# 4.7
content('Phần 4 (7/7) — Đa luồng & đồng bộ Mutex', [
    ('Vấn đề: nhiều luồng truy cập 1 chân GPIO → lỗi "Device busy"', 0),
    ('gpio.c: 1 luồng ghi (writer) + 1 luồng đọc (reader)', 0),
    ('Đồng bộ bằng pthread_mutex → truy cập loại trừ tương hỗ', 0),
    ('Liên kết -lgpiod -lpthread khi biên dịch chéo', 0),
    ('Kết quả: 10 sự kiện Pin=1/0 xen kẽ, không tranh chấp', 0),
    ('→ Chứng minh Mutex phân chia tài nguyên hiệu quả', 1),
], img='Chụp terminal QEMU chạy ./gpio in 10 dòng "Event n: Pin=.. at t=.."')

# ===================== KET LUAN =====================
content('Kết luận', [
    ('Hoàn thành đủ 4 phần của đề tài', 0),
    ('Làm chủ: lập trình shell, C system call, module nhân, HĐH nhúng', 0),
    ('Hiểu cơ chế tương tác giữa chương trình và nhân Linux', 0),
    ('Phần 4: dựng & chạy ARM Linux trên QEMU, điều khiển GPIO + đa luồng', 0),
    ('Hướng phát triển:', 0),
    ('Mở rộng module thành driver thiết bị thực tế', 1),
    ('Bổ sung giao diện cho các chương trình quản lý', 1),
])

# ===================== CAM ON =====================
s = prs.slides.add_slide(blank)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.9),
                         prs.slide_width, Inches(1.8))
bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
tb = s.shapes.add_textbox(Inches(0.8), Inches(3.15), Inches(11.7), Inches(1.4))
p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = 'CẢM ƠN THẦY/CÔ VÀ CÁC BẠN ĐÃ LẮNG NGHE!'
r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = WHITE

prs.save(r'v:\Project\ltnln\Slide_LapTrinhNhanLinux.pptx')
print('Da tao Slide_LapTrinhNhanLinux.pptx — so slide:', len(prs.slides._sldIdLst))
